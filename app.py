import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import io
from datetime import datetime
from uuid import uuid4
from PIL import Image

# --------------------------------------------------
# Page setup
# --------------------------------------------------
st.set_page_config(page_title="Field Inspection Report Generator", layout="wide")
st.title("Field Inspection Report Generator")
st.markdown("---")

# --------------------------------------------------
# Session state
# --------------------------------------------------
if "report_items" not in st.session_state:
    st.session_state.report_items = []
if "generated_ppt_binary" not in st.session_state:
    st.session_state.generated_ppt_binary = None
if "generated_filename" not in st.session_state:
    st.session_state.generated_filename = ""
if "uploader_id" not in st.session_state:
    st.session_state.uploader_id = 0
if "debug_log" not in st.session_state:
    st.session_state.debug_log = []

# Ensure stable IDs for all items
for item in st.session_state.report_items:
    if "id" not in item:
        item["id"] = uuid4().hex

# --------------------------------------------------
# Helpers
# --------------------------------------------------
def log(msg: str):
    ts = datetime.now().strftime("%H:%M:%S")
    st.session_state.debug_log.append(f"[{ts}] {msg}")

def get_image_wh(uploaded_file):
    """Return (w, h) and reset pointer so ppt add_picture still works."""
    try:
        uploaded_file.seek(0)
    except Exception:
        pass
    img = Image.open(uploaded_file)
    w, h = img.size
    try:
        uploaded_file.seek(0)
    except Exception:
        pass
    return w, h

def add_border(slide, x, y, w, h, rgb=RGBColor(0, 0, 0), width_pt=1):
    """Reliable border for pictures: draw transparent rectangle over image."""
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    border.fill.background()  # transparent fill
    border.line.color.rgb = rgb
    border.line.width = Pt(width_pt)
    return border

# --------------------------------------------------
# Callbacks
# --------------------------------------------------
def add_entry_callback():
    uploader_key = f"uploader_{st.session_state.uploader_id}"
    uploaded_file = st.session_state.get(uploader_key)
    description = (st.session_state.get("entry_desc") or "").strip()
    category = st.session_state.get("cat_selector", "Exterior")
    custom_cat = (st.session_state.get("custom_cat_input") or "").strip()

    final_cat = category
    if category == "Other..." and custom_cat:
        final_cat = custom_cat
    elif category == "Other..." and not custom_cat:
        final_cat = "Other"

    if uploaded_file and description:
        st.session_state.report_items.append({
            "id": uuid4().hex,
            "category": final_cat,
            "text": description,
            "image": uploaded_file
        })
        st.session_state["entry_desc"] = ""
        st.session_state.uploader_id += 1
        st.session_state.generated_ppt_binary = None
    else:
        st.error("Please provide both an image and a description.")

def delete_item_callback(index):
    st.session_state.report_items.pop(index)
    st.session_state.generated_ppt_binary = None

def update_item_text(item_id):
    for it in st.session_state.report_items:
        if it["id"] == item_id:
            it["text"] = (st.session_state.get(f"desc_{item_id}") or "").strip()
            break
    st.session_state.generated_ppt_binary = None

def update_item_category(item_id):
    selected = st.session_state.get(f"cat_sel_{item_id}", "Exterior")
    custom = (st.session_state.get(f"cat_other_{item_id}") or "").strip()

    final_cat = selected
    if selected == "Other..." and custom:
        final_cat = custom
    elif selected == "Other..." and not custom:
        final_cat = "Other"

    for it in st.session_state.report_items:
        if it["id"] == item_id:
            it["category"] = final_cat
            break
    st.session_state.generated_ppt_binary = None

def update_item_image(item_id):
    uploaded = st.session_state.get(f"img_{item_id}")
    if uploaded:
        for it in st.session_state.report_items:
            if it["id"] == item_id:
                it["image"] = uploaded
                break
        st.session_state.generated_ppt_binary = None

def move_item(from_index, to_index):
    items = st.session_state.report_items
    if from_index < 0 or from_index >= len(items):
        return
    if to_index < 0 or to_index >= len(items):
        return
    item = items.pop(from_index)
    items.insert(to_index, item)
    st.session_state.generated_ppt_binary = None

def move_up(i):
    if i > 0:
        move_item(i, i - 1)

def move_down(i):
    if i < len(st.session_state.report_items) - 1:
        move_item(i, i + 1)

def move_top(i):
    if i > 0:
        move_item(i, 0)

def move_bottom(i):
    last = len(st.session_state.report_items) - 1
    if i < last:
        move_item(i, last)

# --------------------------------------------------
# Sidebar
# --------------------------------------------------
with st.sidebar:
    st.header("Report Settings")
    report_title = st.text_input("Report Title", "Field Inspection Report")
    date_option = st.selectbox(
        "Date Format",
        ["Month & Year", "Date Only (MM-DD-YYYY)", "Date & Time", "Custom Text"]
    )

    report_subtitle = ""
    filename_suffix = ""

    if date_option == "Custom Text":
        report_subtitle = st.text_input("Subtitle Text", "January 2026")
        filename_suffix = report_subtitle.replace(" ", "_").replace("/", "-")
    else:
        selected_date = st.date_input("Select Date", datetime.now())
        if date_option == "Month & Year":
            report_subtitle = selected_date.strftime("%B %Y")
            filename_suffix = selected_date.strftime("%b_%Y")
        elif date_option == "Date Only (MM-DD-YYYY)":
            report_subtitle = selected_date.strftime("%m-%d-%Y")
            filename_suffix = selected_date.strftime("%m-%d-%Y")
        else:
            selected_time = st.time_input("Select Time", datetime.now().time())
            final_dt = datetime.combine(selected_date, selected_time)
            report_subtitle = final_dt.strftime("%m-%d-%Y %H:%M")
            filename_suffix = final_dt.strftime("%m-%d-%Y_%H%M")

    st.divider()
    st.caption("**Preview:**")
    st.info(report_subtitle)

    clean_title = report_title.replace(" ", "_")
    final_filename = f"{clean_title}_{filename_suffix}.pptx"
    st.caption(f"**Filename:** {final_filename}")

    st.divider()
    debug_mode = st.checkbox(
        "Debug mode",
        value=False,
        help="Shows PPT build logs for troubleshooting (safe for normal users to ignore)."
    )

# --------------------------------------------------
# Batch upload
# --------------------------------------------------
with st.expander("Batch Upload (Add Multiple Images)", expanded=False):
    st.write("Select all images in your folder and drag them here.")
    batch_files = st.file_uploader(
        "Select Multiple Images",
        type=["png", "jpg", "jpeg"],
        accept_multiple_files=True
    )

    if st.button("Add All Batch Images", type="primary"):
        if batch_files:
            for f in batch_files:
                st.session_state.report_items.append({
                    "id": uuid4().hex,
                    "category": "Exterior",
                    "text": "",
                    "image": f
                })
            st.session_state.generated_ppt_binary = None
            st.success(f"Added {len(batch_files)} images! Scroll down to edit.")
        else:
            st.warning("No files selected.")

# --------------------------------------------------
# Quick add
# --------------------------------------------------
st.subheader("Add Single Entry")
c1, c2 = st.columns([1, 1])

with c1:
    standard_options = ["Exterior", "Interior"]
    st.selectbox("Category", standard_options + ["Other..."], key="cat_selector")
    if st.session_state.get("cat_selector") == "Other...":
        st.text_input("Enter Custom Category", key="custom_cat_input")
    st.text_area("Description", height=150, placeholder="Enter observation here...", key="entry_desc")

with c2:
    dynamic_key = f"uploader_{st.session_state.uploader_id}"
    st.file_uploader("Upload Image (Single)", type=["png", "jpg", "jpeg"], key=dynamic_key)

st.button("Add Entry", type="primary", on_click=add_entry_callback)

# --------------------------------------------------
# Current entries
# --------------------------------------------------
if st.session_state.report_items:
    st.markdown("---")
    st.subheader(f"Current Entries ({len(st.session_state.report_items)})")
    st.caption("Shown in page order (top = Page 1). Reorder with arrows. Edit everything inline.")

    for i, item in enumerate(st.session_state.report_items):
        item_id = item["id"]

        st.markdown(
            f"""
            <div style="
                padding: 14px 16px;
                border-radius: 12px;
                border: 1px solid #e6e6e6;
                background: #f7f8fa;
                margin-bottom: 10px;
            ">
                <div style="display:flex; justify-content:space-between; align-items:center;">
                    <div style="font-weight:800; font-size:16px; color:#111111;">
                        Page {i+1}
                    </div>
                    <div style="
                        background:#111827;
                        color:#ffffff;
                        padding:4px 10px;
                        border-radius:999px;
                        font-size:12px;
                        font-weight:800;
                    ">
                        {item["category"]}
                    </div>
                </div>
            </div>
            """,
            unsafe_allow_html=True
        )

        col_img, col_fields, col_actions = st.columns([2, 6, 2])

        with col_img:
            st.image(item["image"], use_container_width=True)
            st.file_uploader(
                "Replace image",
                type=["png", "jpg", "jpeg"],
                key=f"img_{item_id}",
                on_change=update_item_image,
                args=(item_id,),
                label_visibility="collapsed",
            )
            st.caption("Drop a new image to replace")

        with col_fields:
            base = ["Exterior", "Interior"]
            cur_cat = item["category"]
            is_std = cur_cat in base
            default_ix = base.index(cur_cat) if is_std else len(base)

            st.selectbox(
                "Category",
                base + ["Other..."],
                index=default_ix,
                key=f"cat_sel_{item_id}",
                on_change=update_item_category,
                args=(item_id,),
            )

            if st.session_state.get(f"cat_sel_{item_id}") == "Other...":
                st.text_input(
                    "Custom category",
                    value="" if is_std else cur_cat,
                    key=f"cat_other_{item_id}",
                    on_change=update_item_category,
                    args=(item_id,),
                )

            st.text_area(
                "Description",
                value=item.get("text", ""),
                height=140,
                key=f"desc_{item_id}",
                on_change=update_item_text,
                args=(item_id,),
                placeholder="Type the observation here...",
            )

        with col_actions:
            st.button("Top", key=f"top_{item_id}", on_click=move_top, args=(i,), use_container_width=True, disabled=(i == 0))
            st.button("Up", key=f"up_{item_id}", on_click=move_up, args=(i,), use_container_width=True, disabled=(i == 0))
            st.button("Down", key=f"down_{item_id}", on_click=move_down, args=(i,), use_container_width=True, disabled=(i == len(st.session_state.report_items) - 1))
            st.button("Bottom", key=f"bottom_{item_id}", on_click=move_bottom, args=(i,), use_container_width=True, disabled=(i == len(st.session_state.report_items) - 1))
            st.divider()
            st.button("Delete", key=f"delete_{item_id}", on_click=delete_item_callback, args=(i,), use_container_width=True)

        st.divider()

# --------------------------------------------------
# Debug log panel
# --------------------------------------------------
if debug_mode:
    with st.expander("Debug Log", expanded=False):
        st.code("\n".join(st.session_state.debug_log) or "No logs yet.")
        st.download_button(
            "Download log (.txt)",
            data="\n".join(st.session_state.debug_log),
            file_name="debug_log.txt",
            mime="text/plain",
            use_container_width=True
        )

# --------------------------------------------------
# Generate PPT
# --------------------------------------------------
if st.session_state.report_items:
    if st.session_state.generated_ppt_binary is None:
        if st.button("Generate Report", type="primary", use_container_width=True):
            st.session_state.debug_log = []
            log("Starting PPT generation...")

            prs = Presentation()

            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = report_title
            slide.placeholders[1].text = report_subtitle

            # Constants (footer position SAME for portrait & landscape)
            SLIDE_W = Inches(10)
            SLIDE_H = Inches(7.5)
            M = Inches(0.5)

            # Keep footer at the same Y for both layouts
            FOOTER_H = Inches(0.50)                 # same feel as your vertical layout
            FOOTER_Y = SLIDE_H - FOOTER_H           # puts top of footer at 7.0"
            CONTENT_BOTTOM = FOOTER_Y - Inches(0.15)  # content must stay above this

            header_color = RGBColor(176, 196, 222)
            border_color = RGBColor(0, 0, 0)

            for index, item in enumerate(st.session_state.report_items):
                slide = prs.slides.add_slide(prs.slide_layouts[6])

                # Background
                bg = slide.background
                bg.fill.solid()
                bg.fill.fore_color.rgb = RGBColor(200, 210, 215)

                # Detect orientation via ratio
                try:
                    w, h = get_image_wh(item["image"])
                    ratio = (w / h) if h else 1.0
                except Exception as e:
                    w, h, ratio = 0, 0, 1.0
                    log(f"Page {index+1}: ERROR reading image size -> {e}")

                is_landscape = ratio >= 1.10
                log(f"Page {index+1}: image {w}x{h}, ratio={ratio:.2f}, landscape={is_landscape}")

                if not is_landscape:
                    # ===== Portrait layout: header + desc left, image right =====
                    TOP_Y = Inches(0.7)
                    GAP = Inches(0.2)
                    COL = Inches(4.4)
                    HEAD = Inches(0.8)
                    BODY = Inches(5.4)
                    IMG_H = HEAD + BODY

                    # Header (left)
                    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, TOP_Y, COL, HEAD)
                    header.fill.solid()
                    header.fill.fore_color.rgb = header_color
                    header.line.color.rgb = border_color
                    header.text = item["category"]
                    header.text_frame.margin_left = Inches(0.2)
                    header.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    p = header.text_frame.paragraphs[0]
                    p.font.bold = True
                    p.font.size = Pt(26)
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.alignment = PP_ALIGN.LEFT

                    # Description (left)
                    desc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, TOP_Y + HEAD, COL, BODY)
                    desc.fill.solid()
                    desc.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    desc.line.color.rgb = border_color
                    tf = desc.text_frame
                    tf.clear()
                    tf.text = item.get("text", "")
                    tf.word_wrap = True
                    tf.margin_left = Inches(0.2)
                    tf.margin_top = Inches(0.2)
                    tf.vertical_anchor = MSO_ANCHOR.TOP
                    p = tf.paragraphs[0]
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.alignment = PP_ALIGN.LEFT

                    # Image (right)
                    img_x = M + COL + GAP
                    try:
                        item["image"].seek(0)
                    except Exception:
                        pass
                    slide.shapes.add_picture(item["image"], img_x, TOP_Y, width=COL, height=IMG_H)
                    add_border(slide, img_x, TOP_Y, COL, IMG_H, rgb=border_color, width_pt=1)

                else:
                    # ===== Landscape layout (requested):
                    # Category + Description together (top), image BELOW them, footer stays same spot =====
                
                    TOP_Y = Inches(0.7)          # same as portrait
                    FULL_W = SLIDE_W - (M * 2)   # full width usable
                    GAP = Inches(0.2)
                
                    # Make Category + Description look like your portrait left stack
                    HEAD = Inches(0.8)           # same header height feel
                    DESC_H = Inches(1.45)        # tune this if you want more/less text space
                
                    # Header full width, same Y as portrait header
                    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, TOP_Y, FULL_W, HEAD)
                    header.fill.solid()
                    header.fill.fore_color.rgb = header_color
                    header.line.color.rgb = border_color
                    header.text = item["category"]
                    header.text_frame.margin_left = Inches(0.2)
                    header.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    p = header.text_frame.paragraphs[0]
                    p.font.bold = True
                    p.font.size = Pt(26)
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.alignment = PP_ALIGN.LEFT
                
                    # Description directly under header (connected look)
                    desc_y = TOP_Y + HEAD
                    desc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, desc_y, FULL_W, DESC_H)
                    desc.fill.solid()
                    desc.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    desc.line.color.rgb = border_color
                
                    tf = desc.text_frame
                    tf.clear()
                    tf.text = item.get("text", "")
                    tf.word_wrap = True
                    tf.margin_left = Inches(0.2)
                    tf.margin_top = Inches(0.2)
                    tf.vertical_anchor = MSO_ANCHOR.TOP
                    p = tf.paragraphs[0]
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.alignment = PP_ALIGN.LEFT
                
                    # Image goes BELOW description and auto-fits above footer
                    img_y = desc_y + DESC_H + GAP
                    img_h = CONTENT_BOTTOM - img_y  # CONTENT_BOTTOM already protects the footer
                
                    # Safety clamp so it never goes negative
                    if img_h < Inches(2.0):
                        # If someone writes a huge description, reduce desc height to protect the image
                        DESC_H = Inches(1.0)
                        # Move desc shape height down
                        desc.height = DESC_H
                        img_y = desc_y + DESC_H + GAP
                        img_h = CONTENT_BOTTOM - img_y
                
                    try:
                        item["image"].seek(0)
                    except Exception:
                        pass
                
                    slide.shapes.add_picture(item["image"], M, img_y, width=FULL_W, height=img_h)
                    add_border(slide, M, img_y, FULL_W, img_h, rgb=border_color, width_pt=1)


                # Footer (same placement for BOTH layouts)
                footer_box = slide.shapes.add_textbox(M, FOOTER_Y, Inches(6), FOOTER_H)
                fp = footer_box.text_frame.paragraphs[0]
                fp.text = report_title
                fp.font.size = Pt(10)
                fp.font.color.rgb = RGBColor(80, 80, 80)

                page_box = slide.shapes.add_textbox(SLIDE_W - M - Inches(2), FOOTER_Y, Inches(2), FOOTER_H)
                pp = page_box.text_frame.paragraphs[0]
                pp.text = f"Page {index + 1}"
                pp.font.size = Pt(10)
                pp.font.color.rgb = RGBColor(80, 80, 80)
                pp.alignment = PP_ALIGN.RIGHT

            # Save once
            binary = io.BytesIO()
            prs.save(binary)
            binary.seek(0)

            st.session_state.generated_ppt_binary = binary
            st.session_state.generated_filename = final_filename
            log("PPT generation complete.")
            st.rerun()

    else:
        st.download_button(
            label=f"Download {st.session_state.generated_filename}",
            data=st.session_state.generated_ppt_binary,
            file_name=st.session_state.generated_filename,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary",
            use_container_width=True,
        )

        if st.button("Reset / Start New Report", use_container_width=True):
            st.session_state.report_items = []
            st.session_state.generated_ppt_binary = None
            st.session_state.uploader_id += 1
            st.rerun()
